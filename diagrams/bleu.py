import os
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches

# -----------------------------
# Create output directory
# -----------------------------
os.makedirs("diagrams", exist_ok=True)

# -----------------------------
# COMET Scores
# -----------------------------

data = {
    "Flores+":      [1.49, 1.40, 0.85, 0.72, 0.87, 13.41, 14.34],
    "IN22-Conv":    [2.04, 1.99, 1.42, 1.25, 1.29, 9.73, 11.85],
    "IN22-Gen":     [2.42, 2.40, 1.93, 2.01, 1.96, 11.44, 12.84],
    "NIOS":         [1.42, 1.48, 1.02, 1.09, 1.12, 3.23, 4.47]
}

models = [
    "ZS",
    "R (sa-en)",
    "R (sa-en-hi)",
    "R (sa-hi)",
    "R (sa-hi-en)",
    "NLLB 1.3B",
    "NLLB 3.3B"
]

# -----------------------------
# Create DataFrame
# -----------------------------
df = pd.DataFrame(data, index=models)

# -----------------------------
# Plot Heatmap
# -----------------------------

plt.figure(figsize=(12, 7))

ax = sns.heatmap(
    df,
    annot=True,
    fmt=".2f",
    cmap="YlGnBu",
    linewidths=0.5,
    cbar=True,
    annot_kws={
        "size": 20,
        "weight": "normal"   # Normal annotation text
    },
    cbar_kws={
        "shrink": 0.9
    }
)

# -----------------------------
# Font settings
# -----------------------------

# Axis labels
ax.set_xlabel("Datasets", fontsize=20, fontweight='normal')
ax.set_ylabel("Models", fontsize=20, fontweight='normal')

# Tick labels
ax.set_xticklabels(
    ax.get_xticklabels(),
    rotation=0,
    fontsize=20,
    fontweight='normal'
)

ax.set_yticklabels(
    ax.get_yticklabels(),
    rotation=0,
    fontsize=20,
    fontweight='normal'
)

# Colorbar font size
cbar = ax.collections[0].colorbar
cbar.ax.tick_params(labelsize=18)

# Remove title
plt.title("")

# Tight layout
plt.tight_layout()

# -----------------------------
# Save figure
# -----------------------------

png_path = "diagrams/bleu.png"
pdf_path = "diagrams/bleu.pdf"
pptx_path = "diagrams/bleu.pptx"

plt.savefig(png_path, dpi=300, bbox_inches='tight')
plt.savefig(pdf_path, bbox_inches='tight')

# Show plot
plt.show()

# -----------------------------
# Save to PowerPoint (.pptx)
# -----------------------------

prs = Presentation()

# Blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add image to slide
slide.shapes.add_picture(
    png_path,
    Inches(0.5),
    Inches(0.5),
    width=Inches(12.5)
)

# Save PPTX
prs.save(pptx_path)

print("Saved:")
print(f" - {png_path}")
print(f" - {pdf_path}")
print(f" - {pptx_path}")