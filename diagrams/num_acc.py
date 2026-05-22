import os
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt

# -----------------------------
# Create output directory
# -----------------------------
os.makedirs("diagrams", exist_ok=True)

# -----------------------------
# Numeral Accuracy (%)
# -----------------------------

data = {
    "Flores+":   [21.90, 85.71, 89.52],
    "IN22-Conv": [83.54, 72.15, 74.68],
    "IN22-Gen":  [74.42, 86.05, 89.37],
    "NIOS":      [86.00, 77.33, 84.00]
}

models = [
    "ZS",
    "NLLB 1.3B",
    "NLLB 3.3B"
]

# Create DataFrame
df = pd.DataFrame(data, index=models)

# -----------------------------
# Plot Heatmap
# -----------------------------

fig, ax = plt.subplots(figsize=(8, 4.8))

sns.heatmap(
    df,
    annot=True,
    fmt=".2f",
    cmap="YlOrRd",
    linewidths=0.5,
    cbar=True,
    square=False,
    annot_kws={
        "fontsize": 20,
        "fontweight": "normal"
    },
    cbar_kws={
        "shrink": 0.9,
        "pad": 0.02
    },
    ax=ax
)

# Remove title
plt.title("")

# Axis labels
plt.xlabel("Datasets", fontsize=20)
plt.ylabel("Models", fontsize=20)

# -----------------------------
# Adjust Dataset Names Properly
# -----------------------------

ax.set_xticklabels(
    ax.get_xticklabels(),
    rotation=30,
    ha='right',              # keeps labels ending before heatmap
    rotation_mode='anchor',
    fontsize=20
)

ax.set_yticklabels(
    ax.get_yticklabels(),
    rotation=0,
    fontsize=20
)

# Add padding between labels and heatmap
ax.tick_params(axis='x', pad=10)

# -----------------------------
# Colorbar Settings
# -----------------------------

cbar = ax.collections[0].colorbar
cbar.ax.tick_params(labelsize=20)

# Remove colorbar label
cbar.set_label("")

# -----------------------------
# Remove top and right ticks
# -----------------------------

ax.tick_params(
    top=False,
    bottom=True,
    left=True,
    right=False,
    labeltop=False,
    labelbottom=True
)

# -----------------------------
# Extra Space for Tick Labels
# -----------------------------

plt.subplots_adjust(
    bottom=0.32,   # more bottom margin
    left=0.30,
    right=0.95,
    top=0.95
)

# -----------------------------
# Save Outputs
# -----------------------------

png_path = "diagrams/num_acc.png"
pdf_path = "diagrams/num_acc.pdf"

plt.savefig(
    png_path,
    dpi=300,
    bbox_inches='tight'
)

plt.savefig(
    pdf_path,
    bbox_inches='tight'
)

plt.close()

# -----------------------------
# Save in PPTX
# -----------------------------

ppt = Presentation()

# Use blank slide layout
slide_layout = ppt.slide_layouts[6]
slide = ppt.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(
    Inches(0.5),
    Inches(0.2),
    Inches(8),
    Inches(0.5)
)

title_para = title_box.text_frame.paragraphs[0]
title_para.text = "Numeral Accuracy Heatmap"
title_para.font.size = Pt(20)
title_para.font.bold = False

# Add image
slide.shapes.add_picture(
    png_path,
    Inches(0.5),
    Inches(0.8),
    width=Inches(8.5)
)

pptx_path = "diagrams/num_acc.pptx"
ppt.save(pptx_path)

# -----------------------------
# Print Saved Files
# -----------------------------

print("Saved:")
print(f" - {png_path}")
print(f" - {pdf_path}")
print(f" - {pptx_path}")