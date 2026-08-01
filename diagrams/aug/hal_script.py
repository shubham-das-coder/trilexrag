import json
import os
import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt

JSONL_FILE_PATH = "additional/script/total_percentage.jsonl"

OUTPUT_DIR = "diagrams/aug/hal-script"
os.makedirs(OUTPUT_DIR, exist_ok=True)
OUTPUT_BASENAME = "script"

MODEL_NAME_MAP = {
    "qwen3_4b_instruct": "Qwen3 4B",
    "gpt_oss_20b": "GPT-OSS 20B",
    "indictrans2_320m": "IndicTrans2 320M",
    "indictrans2_1b": "IndicTrans2 1B",
    "nllb200_1p3b": "NLLB 1.3B",
    "nllb200_3p3b": "NLLB 3.3B",
}

DATASET_NAME_MAP = {
    "flores_plus": "Flores+",
    "in22_conv": "IN22-Conv",
    "in22_gen": "IN22-Gen",
    "nios": "NIOS",
}

parsed_data = {}

if not os.path.exists(JSONL_FILE_PATH):
    raise FileNotFoundError(f"Input file not found at: {JSONL_FILE_PATH}")

with open(JSONL_FILE_PATH, "r", encoding="utf-8") as f:
    for line in f:
        line = line.strip()
        if not line:
            continue

        row = json.loads(line)

        raw_model = row["model"]
        raw_dataset = row["file"]

        model_name = MODEL_NAME_MAP.get(raw_model, raw_model)
        dataset_name = DATASET_NAME_MAP.get(raw_dataset, raw_dataset)

        if "hallucinated_rows" in row and "dataset_total" in row:
            hallucinated = float(row["hallucinated_rows"])
            total = float(row["dataset_total"])
            perc_val = (hallucinated / total) * 100.0 if total > 0 else 0.0
        else:
            perc_val = float(row.get("percentage", 0.0))

        parsed_data.setdefault(model_name, {})[dataset_name] = perc_val

# Create DataFrame from parsed dictionary
df = pd.DataFrame.from_dict(parsed_data, orient="index")

# Order and align rows (models)
ordered_models = [
    MODEL_NAME_MAP[k]
    for k in MODEL_NAME_MAP
    if MODEL_NAME_MAP[k] in df.index
]

# Order and align columns (datasets)
ordered_datasets = [
    DATASET_NAME_MAP[k]
    for k in DATASET_NAME_MAP
    if DATASET_NAME_MAP[k] in df.columns or DATASET_NAME_MAP[k] in [val for sub in parsed_data.values() for val in sub]
]

# Reindex dataframe to guarantee standard row & column layout
df = df.reindex(index=ordered_models, columns=ordered_datasets)

# Fill any missing values (NaN) with 0.0 if a model has no entries for a dataset
df = df.fillna(0.0)

fig, ax = plt.subplots(figsize=(8, 4.8))

sns.heatmap(
    df,
    annot=True,
    fmt=".2f",
    cmap="YlOrRd",
    linewidths=0.5,
    cbar=True,
    square=False,
    annot_kws={"fontsize": 16, "fontweight": "normal"},
    cbar_kws={"shrink": 0.9, "pad": 0.02},
    ax=ax,
)

plt.title("")
plt.xlabel("Datasets", fontsize=18)
plt.ylabel("Methods", fontsize=18)

ax.set_xticklabels(
    ax.get_xticklabels(),
    rotation=30,
    ha="right",
    rotation_mode="anchor",
    fontsize=16,
)

ax.set_yticklabels(ax.get_yticklabels(), rotation=0, fontsize=16)

ax.tick_params(axis="x", pad=10)

cbar = ax.collections[0].colorbar
cbar.ax.tick_params(labelsize=16)
cbar.set_label("% Hallucinated", fontsize=14)

ax.tick_params(
    top=False,
    bottom=True,
    left=True,
    right=False,
    labeltop=False,
    labelbottom=True,
)

plt.subplots_adjust(bottom=0.32, left=0.30, right=0.95, top=0.95)

png_path = os.path.join(OUTPUT_DIR, f"{OUTPUT_BASENAME}.png")
pdf_path = os.path.join(OUTPUT_DIR, f"{OUTPUT_BASENAME}.pdf")

plt.savefig(png_path, dpi=300, bbox_inches="tight")
plt.savefig(pdf_path, bbox_inches="tight")
plt.close()

ppt = Presentation()
slide_layout = ppt.slide_layouts[6]
slide = ppt.slides.add_slide(slide_layout)

title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.2), Inches(8), Inches(0.5)
)
title_para = title_box.text_frame.paragraphs[0]
title_para.text = "Script Hallucination Heatmap (%)"
title_para.font.size = Pt(20)
title_para.font.bold = False

slide.shapes.add_picture(png_path, Inches(0.5), Inches(0.8), width=Inches(8.5))

pptx_path = os.path.join(OUTPUT_DIR, f"{OUTPUT_BASENAME}.pptx")
ppt.save(pptx_path)

print(f"Process finished successfully! Files saved in: '{OUTPUT_DIR}'")
print(f" - Image PNG:   {png_path}")
print(f" - Vector PDF:  {pdf_path}")
print(f" - PowerPoint:  {pptx_path}")