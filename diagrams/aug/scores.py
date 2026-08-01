import json
import os
import glob
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches

METHOD_FILE_MAP = {
    "Qwen3 ZS": "combined/zs/qwen3_4b_instruct_scores.jsonl",
    "Qwen3 DGT (sa-en)": "combined/se/qwen3_4b_instruct_scores.jsonl",
    "GPT-OSS OS": "combined/os/gpt_oss_20b_scores.jsonl",
    "GPT-OSS FS": "combined/fs/gpt_oss_20b_scores.jsonl",
    "IndicTrans2 320M": "combined/nmt/indictrans2_320m_scores.jsonl",
    "IndicTrans2 1B": "combined/nmt/indictrans2_1b_scores.jsonl",
    "NLLB-1.3B": "combined/nmt/nllb200_1p3b_scores.jsonl",
    "NLLB-3.3B": "combined/nmt/nllb200_3p3b_scores.jsonl"
}

DATASET_NAME_MAP = {
    "flores_plus.jsonl": "Flores+",
    "in22_conv.jsonl": "IN22-Conv",
    "in22_gen.jsonl": "IN22-Gen",
    "nios.jsonl": "NIOS"
}

DATASET_ORDER = ["Flores+", "IN22-Conv", "IN22-Gen", "NIOS"]

METRICS_TO_PLOT = {
    "bleu": {
        "key": "bleu",
        "output_prefix": "bleu",
        "scale_100": True
    },
    "meteor": {
        "key": "meteor",
        "output_prefix": "meteor",
        "scale_100": True
    },
    "comet": {
        "key": "wmt22_comet_da",
        "output_prefix": "comet",
        "scale_100": True
    }
}

X_AXIS_LABEL = "Datasets"
Y_AXIS_LABEL = "Methods"
OUTPUT_DIR = "diagrams/aug/scores"


def load_and_parse_data(metric_cfg):
    metric_key = metric_cfg["key"]
    should_scale = metric_cfg["scale_100"]

    matrix_data = {}

    for method_name, file_path in METHOD_FILE_MAP.items():
        method_scores = {}

        matching_files = glob.glob(file_path) if "*" in file_path else [file_path]

        for path in matching_files:
            if not os.path.exists(path):
                print(f"Warning: File not found - {path}")
                continue

            with open(path, "r", encoding="utf-8") as f:
                for line in f:
                    if not line.strip():
                        continue

                    entry = json.loads(line)

                    raw_file_name = entry.get("file", "")

                    dataset_name = DATASET_NAME_MAP.get(
                        raw_file_name,
                        raw_file_name.replace(".jsonl", "").replace("_", " ").title()
                    )

                    val = entry.get(metric_key, 0.0)

                    if should_scale and val <= 1.0 and val > 0.0:
                        val *= 100.0

                    method_scores[dataset_name] = val

        matrix_data[method_name] = method_scores

    df = pd.DataFrame.from_dict(matrix_data, orient="index")
    df = df.reindex([m for m in METHOD_FILE_MAP.keys() if m in df.index])

    existing_cols = [col for col in DATASET_ORDER if col in df.columns]
    remaining_cols = [col for col in df.columns if col not in existing_cols]
    df = df[existing_cols + remaining_cols]

    return df


os.makedirs(OUTPUT_DIR, exist_ok=True)

for metric_name, cfg in METRICS_TO_PLOT.items():
    print(f"\nProcessing {metric_name.upper()}...")

    df = load_and_parse_data(cfg)

    if df.empty:
        print(f"No data found for {metric_name}. Skipping...")
        continue

    fig_width = max(10, len(df.columns) * 1.50)
    fig_height = max(4, len(df.index) * 0.75)

    plt.figure(figsize=(fig_width, fig_height))

    ax = sns.heatmap(
        df,
        annot=True,
        fmt=".2f",
        cmap="YlGnBu",
        linewidths=0.5,
        cbar=True,
        square=False,
        annot_kws={
            "size": 18,
            "weight": "normal"
        },
        cbar_kws={
            "shrink": 0.85
        }
    )

    ax.set_xlabel(X_AXIS_LABEL, fontsize=18, fontweight="normal", labelpad=12)
    ax.set_ylabel(Y_AXIS_LABEL, fontsize=18, fontweight="normal", labelpad=12)

    ax.set_xticklabels(
        ax.get_xticklabels(),
        rotation=0,
        fontsize=16,
        fontweight="normal"
    )

    ax.set_yticklabels(
        ax.get_yticklabels(),
        rotation=0,
        fontsize=16,
        fontweight="normal"
    )

    cbar = ax.collections[0].colorbar
    cbar.ax.tick_params(labelsize=14)

    plt.title("")
    plt.tight_layout()

    prefix = cfg["output_prefix"]
    png_path = os.path.join(OUTPUT_DIR, f"{prefix}.png")
    pdf_path = os.path.join(OUTPUT_DIR, f"{prefix}.pdf")
    pptx_path = os.path.join(OUTPUT_DIR, f"{prefix}.pptx")

    plt.savefig(png_path, dpi=300, bbox_inches="tight")
    plt.savefig(pdf_path, bbox_inches="tight")
    plt.close()

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.add_picture(
        png_path,
        Inches(0.5),
        Inches(0.5),
        width=Inches(12.5)
    )

    prs.save(pptx_path)

    print(" Saved:")
    print(f"  - {png_path}")
    print(f"  - {pdf_path}")
    print(f"  - {pptx_path}")

print("\nDone! All diagrams have been generated.")