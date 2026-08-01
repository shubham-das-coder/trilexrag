import os
import json
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT_DIR = "diagrams/aug/hal-number"  
os.makedirs(OUTPUT_DIR, exist_ok=True)

OUTPUT_BASENAME = "number"

JSONL_FILES = {
    "Qwen3 4B": "additional/num/v2/qwen3_4b_instruct_stats.jsonl",
    "GPT-OSS 20B": "additional/num/v2/gpt_oss_20b_stats.jsonl",
    "IndicTrans2 320M": "additional/num/v2/indictrans2_320m_stats.jsonl",
    "IndicTrans2 1B": "additional/num/v2/indictrans2_1b_stats.jsonl",
    "NLLB 1.3B": "additional/num/v2/nllb200_1p3b_stats.jsonl",
    "NLLB 3.3B": "additional/num/v2/nllb200_3p3b_stats.jsonl",
}

DATASET_NAME_MAP = {
    "flores_plus": "Flores+",
    "in22_conv": "IN22-Conv",
    "in22_gen": "IN22-Gen",
    "nios": "NIOS"
}

parsed_data = {}

for model_name, filepath in JSONL_FILES.items():
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"File not found for '{model_name}': {filepath}")

    parsed_data[model_name] = {}

    with open(filepath, "r", encoding="utf-8") as f:
        for line in f:
            if not line.strip():
                continue
            row = json.loads(line)
            ds_key = row["file"]
            ds_name = DATASET_NAME_MAP.get(ds_key, ds_key)
            
            parsed_data[model_name][ds_name] = float(row["numeral_accuracy"])

df_acc = pd.DataFrame.from_dict(parsed_data, orient="index")

expected_datasets = list(DATASET_NAME_MAP.values())
df_acc = df_acc.reindex(columns=expected_datasets)

df_acc = df_acc.fillna(100.0)

df_error = 100.0 - df_acc

fig, ax = plt.subplots(figsize=(8, 4.8))

sns.heatmap(
    df_error,
    annot=True,
    fmt=".2f",
    cmap="YlOrRd",
    linewidths=0.5,
    cbar=True,
    square=False,
    annot_kws={"fontsize": 20, "fontweight": "normal"},
    cbar_kws={"shrink": 0.9, "pad": 0.02},
    ax=ax
)

plt.title("")
plt.xlabel("Datasets", fontsize=20)
plt.ylabel("Methods", fontsize=20)

ax.set_xticklabels(
    ax.get_xticklabels(),
    rotation=30,
    ha='right',
    rotation_mode='anchor',
    fontsize=20
)

ax.set_yticklabels(
    ax.get_yticklabels(),
    rotation=0,
    fontsize=20
)

ax.tick_params(axis='x', pad=10)

cbar = ax.collections[0].colorbar
cbar.ax.tick_params(labelsize=20)
cbar.set_label("")

ax.tick_params(
    top=False,
    bottom=True,
    left=True,
    right=False,
    labeltop=False,
    labelbottom=True
)

plt.subplots_adjust(
    bottom=0.32,
    left=0.30,
    right=0.95,
    top=0.95
)

png_path = os.path.join(OUTPUT_DIR, f"{OUTPUT_BASENAME}.png")
pdf_path = os.path.join(OUTPUT_DIR, f"{OUTPUT_BASENAME}.pdf")

plt.savefig(png_path, dpi=300, bbox_inches='tight')
plt.savefig(pdf_path, bbox_inches='tight')
plt.close()

ppt = Presentation()
slide_layout = ppt.slide_layouts[6]  
slide = ppt.slides.add_slide(slide_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
title_para = title_box.text_frame.paragraphs[0]
title_para.text = "Numeral Error Rate Heatmap"
title_para.font.size = Pt(20)
title_para.font.bold = False

slide.shapes.add_picture(png_path, Inches(0.5), Inches(0.8), width=Inches(8.5))

pptx_path = os.path.join(OUTPUT_DIR, f"{OUTPUT_BASENAME}.pptx")
ppt.save(pptx_path)

print(f"Process finished successfully! Files saved in: '{OUTPUT_DIR}'")
print(f" - Image PNG:   {png_path}")
print(f" - Vector PDF:  {pdf_path}")
print(f" - PowerPoint:  {pptx_path}")