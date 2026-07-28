import json
import os
import glob
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches

# ==========================================
# CONFIGURATION
# ==========================================

# 1. Path mapping: Assign each Model/Method (Y-axis row) to its JSONL result file path.
MODEL_FILE_MAP = {
    "GPT-OSS ZS": "combined/zs/gpt_oss_20b_scores.jsonl",
    "GPT-OSS FS": "combined/fs/gpt_oss_20b_scores.jsonl",
    "GPT-OSS DGT (sa-en)": "combined/se/gpt_oss_20b_scores.jsonl",
    "GPT-OSS DGT (sa-en-hi)": "combined/seh/gpt_oss_20b_scores.jsonl",
    "GPT-OSS DGT (sa-hi)": "combined/sh/gpt_oss_20b_scores.jsonl",
    "GPT-OSS DGT (sa-hi-en)": "combined/she/gpt_oss_20b_scores.jsonl",    
}

# 2. Dataset Display Name Mapping (maps the internal "file" value to chart X-axis labels)
DATASET_NAME_MAP = {
    "flores_plus.jsonl": "Flores+",
    "in22_conv.jsonl": "IN22-Conv",
    "in22_gen.jsonl": "IN22-Gen",
    "nios.jsonl": "NIOS"
}

# Enforce a specific dataset order across the X-axis
DATASET_ORDER = ["Flores+", "IN22-Conv", "IN22-Gen", "NIOS"]

# 3. Metrics to plot and their visual settings
METRICS_TO_PLOT = {
    "bleu": {
        "key": "bleu",
        "output_prefix": "bleu",
        "scale_100": True  # Scale 0.084 -> 8.41
    },
    "meteor": {
        "key": "meteor",
        "output_prefix": "meteor",
        "scale_100": True  # Scale 0.266 -> 26.63
    },
    "comet": {
        "key": "wmt22_comet_da",
        "output_prefix": "comet",
        "scale_100": True  # Scale 0.643 -> 64.33
    }
}

# 4. Global Axis Labels & Output Folder
X_AXIS_LABEL = "Datasets"
Y_AXIS_LABEL = "Methods"
OUTPUT_DIR = "diagrams/aug/gpt-oss"

# ==========================================
# DATA PROCESSING FUNCTION
# ==========================================
def load_and_parse_data(metric_cfg):
    metric_key = metric_cfg["key"]
    should_scale = metric_cfg["scale_100"]
    
    matrix_data = {}

    for model_name, file_path in MODEL_FILE_MAP.items():
        model_scores = {}
        
        # Expand wildcard if provided, or use exact file path
        matching_files = glob.glob(file_path) if "*" in file_path else [file_path]
        
        for path in matching_files:
            if not os.path.exists(path):
                print(f"Warning: File not found - {path}")
                continue
                
            with open(path, "r", encoding="utf-8") as f:
                for line in f:
                    if not line.strip():
                        continue
                    entry = json.loads(line)
                    
                    # Extract raw filename identifier
                    raw_file_name = entry.get("file", "")
                    
                    # Map to clean dataset label
                    dataset_name = DATASET_NAME_MAP.get(
                        raw_file_name,
                        raw_file_name.replace(".jsonl", "").replace("_", " ").title()
                    )
                    
                    val = entry.get(metric_key, 0.0)
                    
                    # Multiply by 100 if decimal scores need standard formatting
                    if should_scale and val <= 1.0 and val > 0.0:
                        val *= 100.0
                        
                    model_scores[dataset_name] = val
                    
        matrix_data[model_name] = model_scores

    # Construct DataFrame
    df = pd.DataFrame.from_dict(matrix_data, orient='index')
    
    # Reorder columns to match requested dataset sequence
    existing_cols = [col for col in DATASET_ORDER if col in df.columns]
    remaining_cols = [col for col in df.columns if col not in existing_cols]
    df = df[existing_cols + remaining_cols]
    
    return df

# ==========================================
# GENERATION LOOP FOR BLEU, METEOR, COMET
# ==========================================
os.makedirs(OUTPUT_DIR, exist_ok=True)

for metric_name, cfg in METRICS_TO_PLOT.items():
    print(f"\nProcessing {metric_name.upper()}...")
    
    df = load_and_parse_data(cfg)
    
    if df.empty:
        print(f"No data found for {metric_name}. Skipping...")
        continue

    # Plot size based on row/column count
    fig_width = max(10, len(df.columns) * 2.5)
    fig_height = max(6, len(df.index) * 1.0)

    plt.figure(figsize=(fig_width, fig_height))

    ax = sns.heatmap(
        df,
        annot=True,
        fmt=".2f",
        cmap="YlGnBu",
        linewidths=0.5,
        cbar=True,
        annot_kws={
            "size": 18,
            "weight": "normal"
        },
        cbar_kws={
            "shrink": 0.85
        }
    )

    # Axis Labels
    ax.set_xlabel(X_AXIS_LABEL, fontsize=20, fontweight='normal', labelpad=12)
    ax.set_ylabel(Y_AXIS_LABEL, fontsize=20, fontweight='normal', labelpad=12)

    # Ticks Formatting
    ax.set_xticklabels(
        ax.get_xticklabels(),
        rotation=0,
        fontsize=18,
        fontweight='normal'
    )

    ax.set_yticklabels(
        ax.get_yticklabels(),
        rotation=0,
        fontsize=18,
        fontweight='normal'
    )

    # Colorbar Tick Size
    cbar = ax.collections[0].colorbar
    cbar.ax.tick_params(labelsize=16)

    plt.title("")
    plt.tight_layout()

    # Define File Paths
    prefix = cfg["output_prefix"]
    png_path = os.path.join(OUTPUT_DIR, f"{prefix}.png")
    pdf_path = os.path.join(OUTPUT_DIR, f"{prefix}.pdf")
    pptx_path = os.path.join(OUTPUT_DIR, f"{prefix}.pptx")

    # Save PNG and PDF
    plt.savefig(png_path, dpi=300, bbox_inches='tight')
    plt.savefig(pdf_path, bbox_inches='tight')
    plt.close()

    # Save PowerPoint Slide
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.add_picture(
        png_path,
        Inches(0.5),
        Inches(0.5),
        width=Inches(12.5)
    )

    prs.save(pptx_path)

    print(f" Saved:")
    print(f"  - {png_path}")
    print(f"  - {pdf_path}")
    print(f"  - {pptx_path}")

print("\nDone! All diagrams have been generated.")