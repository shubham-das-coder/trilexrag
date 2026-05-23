import os
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches

# -----------------------------
# Create output directory
# -----------------------------
os.makedirs("diagrams", exist_ok=True)

# -----------------------------
# COMET Scores
# -----------------------------

data = {
    "Flores+":      [12.43, 12.60, 9.79, 9.86, 9.80, 33.01, 34.66],
    "IN22-Conv":    [15.62, 15.79, 12.82, 12.88, 12.79, 27.48, 30.30],
    "IN22-Gen":     [14.68, 14.82, 12.50, 12.53, 12.49, 32.08, 34.22],
    "NIOS":         [13.95, 14.91, 12.08, 12.45, 11.96, 15.98, 18.07]
}

models = [
    "ZS",
    "DGT (sa-en)",
    "DGT (sa-en-hi)",
    "DGT (sa-hi)",
    "DGT (sa-hi-en)",
    "NLLB 1.3B",
    "NLLB 3.3B"
]

# -----------------------------
# Create DataFrame
# -----------------------------
df = pd.DataFrame(data, index=models)

# -----------------------------
# Plot Heatmap
# -----------------------------

plt.figure(figsize=(12, 7))

ax = sns.heatmap(
    df,
    annot=True,
    fmt=".2f",
    cmap="YlGnBu",
    linewidths=0.5,
    cbar=True,
    annot_kws={
        "size": 20,
        "weight": "normal"   # Normal annotation text
    },
    cbar_kws={
        "shrink": 0.9
    }
)

# -----------------------------
# Font settings
# -----------------------------

# Axis labels
ax.set_xlabel("Datasets", fontsize=20, fontweight='normal')
ax.set_ylabel("Methods", fontsize=20, fontweight='normal')

# Tick labels
ax.set_xticklabels(
    ax.get_xticklabels(),
    rotation=0,
    fontsize=20,
    fontweight='normal'
)

ax.set_yticklabels(
    ax.get_yticklabels(),
    rotation=0,
    fontsize=20,
    fontweight='normal'
)

# Colorbar font size
cbar = ax.collections[0].colorbar
cbar.ax.tick_params(labelsize=18)

# Remove title
plt.title("")

# Tight layout
plt.tight_layout()

# -----------------------------
# Save figure
# -----------------------------

png_path = "diagrams/meteor.png"
pdf_path = "diagrams/meteor.pdf"
pptx_path = "diagrams/meteor.pptx"

plt.savefig(png_path, dpi=300, bbox_inches='tight')
plt.savefig(pdf_path, bbox_inches='tight')

# Show plot
plt.show()

# -----------------------------
# Save to PowerPoint (.pptx)
# -----------------------------

prs = Presentation()

# Blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add image to slide
slide.shapes.add_picture(
    png_path,
    Inches(0.5),
    Inches(0.5),
    width=Inches(12.5)
)

# Save PPTX
prs.save(pptx_path)

print("Saved:")
print(f" - {png_path}")
print(f" - {pdf_path}")
print(f" - {pptx_path}")