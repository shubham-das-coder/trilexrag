import json
from collections import Counter, defaultdict
from pathlib import Path
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches

folder_path = "additional/script"
output_path = "additional/script/cardinality.jsonl"
plot_path = "additional/script/cardinality.pdf"
png_path = "additional/script/cardinality.png"
ppt_path = "additional/script/cardinality.pptx"

Path(output_path).parent.mkdir(parents=True, exist_ok=True)
Path(plot_path).parent.mkdir(parents=True, exist_ok=True)

model_order = [
    "nllb200_1p3b",
    "nllb200_3p3b",
    "qwen3_4b_instruct"
]

dataset_order = [
    "flores_plus",
    "in22_conv",
    "in22_gen",
    "nios"
]

all_data = defaultdict(lambda: defaultdict(Counter))

with open(output_path, "w", encoding="utf-8") as out_f:
    for model_name in model_order:
        path = Path(folder_path) / f"{model_name}_rows.jsonl"

        if not path.exists():
            continue

        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                row = json.loads(line)

                dataset = row["file"]
                cardinality = row["num_scripts"]

                all_data[model_name][dataset][cardinality] += 1

        for dataset in dataset_order:
            if dataset not in all_data[model_name]:
                continue

            counter = all_data[model_name][dataset]

            output_row = {
                "model": model_name,
                "dataset": dataset,
                "cardinality_distribution": dict(sorted(counter.items()))
            }

            out_f.write(json.dumps(output_row, ensure_ascii=False) + "\n")

x_labels = []
combined_data = []

for dataset in dataset_order:
    for model in model_order:
        if dataset not in all_data[model]:
            continue

        x_labels.append(f"{model}\n{dataset}")
        combined_data.append(all_data[model][dataset])

cardinalities = sorted({
    card
    for counter in combined_data
    for card in counter.keys()
})

fig, ax = plt.subplots(figsize=(16, 6))

bottom = [0] * len(combined_data)

for cardinality in cardinalities:
    values = [
        counter.get(cardinality, 0)
        for counter in combined_data
    ]

    ax.bar(
        x_labels,
        values,
        bottom=bottom,
        label=f"{cardinality}"
    )

    bottom = [
        b + v
        for b, v in zip(bottom, values)
    ]

ax.set_xlabel("Model / Dataset")
ax.set_ylabel("Count")
ax.set_title("Stacked Script Cardinality Distribution")

ax.legend(title="Script Cardinality")

plt.xticks(rotation=30)
fig.tight_layout()

with PdfPages(plot_path) as pdf:
    pdf.savefig(fig, bbox_inches="tight")

fig.savefig(png_path, bbox_inches="tight", dpi=300)

prs = Presentation()
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide.shapes.add_picture(
    png_path,
    Inches(0.5),
    Inches(0.5),
    width=Inches(12.5)
)

prs.save(ppt_path)

plt.close(fig)